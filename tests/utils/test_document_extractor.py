"""Tests for deeptutor.utils.document_extractor."""

from __future__ import annotations

import base64
import io

from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
import pytest

from deeptutor.utils import document_extractor as document_extractor_module
from deeptutor.utils.document_extractor import (
    MAX_DOC_BYTES,
    MAX_EXTRACTED_CHARS_PER_DOC,
    CorruptDocumentError,
    DocumentTooLargeError,
    EmptyDocumentError,
    UnsupportedDocumentError,
    extract_documents_from_records,
    extract_text_from_bytes,
    extract_text_from_path,
    is_document_extension,
)

# ---------------------------------------------------------------------------
# Fixtures — generate office docs on the fly
# ---------------------------------------------------------------------------


def _make_docx(paragraphs: list[str]) -> bytes:
    doc = DocxDocument()
    for p in paragraphs:
        doc.add_paragraph(p)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_xlsx(sheets: dict[str, list[list[object]]]) -> bytes:
    wb = Workbook()
    default = wb.active
    first = True
    for name, rows in sheets.items():
        ws = default if first else wb.create_sheet()
        ws.title = name
        for row in rows:
            ws.append(row)
        first = False
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx(slides_text: list[list[str]]) -> bytes:
    prs = Presentation()
    for slide_texts in slides_text:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank-ish layout
        for i, text in enumerate(slide_texts):
            tb = slide.shapes.add_textbox(Inches(1), Inches(1 + i * 0.5), Inches(6), Inches(0.5))
            tb.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# is_document_extension
# ---------------------------------------------------------------------------


class TestIsDocumentExtension:
    def test_office(self) -> None:
        assert is_document_extension("foo.pdf")
        assert is_document_extension("foo.DOCX")
        assert is_document_extension("report.xlsx")
        assert is_document_extension("deck.pptx")

    def test_text_and_code(self) -> None:
        # Any extension in FileTypeRouter.TEXT_EXTENSIONS should be supported.
        assert is_document_extension("notes.txt")
        assert is_document_extension("readme.md")
        assert is_document_extension("module.py")
        assert is_document_extension("config.yaml")
        assert is_document_extension("data.json")
        assert is_document_extension("index.html")
        assert is_document_extension("table.csv")

    def test_unsupported(self) -> None:
        assert not is_document_extension("foo.png")
        assert not is_document_extension("foo.zip")
        assert not is_document_extension("foo.exe")
        assert not is_document_extension("foo")
        assert not is_document_extension("")


# ---------------------------------------------------------------------------
# extract_text_from_bytes — happy paths
# ---------------------------------------------------------------------------


class TestExtractDocx:
    def test_basic_paragraphs(self) -> None:
        data = _make_docx(["Hello world", "Second paragraph", ""])
        text = extract_text_from_bytes("doc.docx", data)
        assert "Hello world" in text
        assert "Second paragraph" in text

    def test_path_helper_can_disable_chat_truncation(self, tmp_path) -> None:
        data = _make_docx(["a" * (MAX_EXTRACTED_CHARS_PER_DOC + 10)])
        path = tmp_path / "long.docx"
        path.write_bytes(data)

        text = extract_text_from_path(path, max_chars=None)

        assert len(text) > MAX_EXTRACTED_CHARS_PER_DOC
        assert "truncated" not in text

    def test_ooxml_fallback_without_python_docx(self, monkeypatch: pytest.MonkeyPatch) -> None:
        monkeypatch.setattr(document_extractor_module, "DocxDocument", None)
        data = _make_docx(["Fallback paragraph", "第二段"])

        text = extract_text_from_bytes("doc.docx", data)

        assert "Fallback paragraph" in text
        assert "第二段" in text


class TestExtractXlsx:
    def test_multiple_sheets(self) -> None:
        data = _make_xlsx(
            {
                "Alpha": [["a1", "b1"], ["a2", 42]],
                "Beta": [["x", "y"]],
            }
        )
        text = extract_text_from_bytes("book.xlsx", data)
        assert "--- Sheet: Alpha ---" in text
        assert "--- Sheet: Beta ---" in text
        assert "a1" in text and "42" in text
        assert "x" in text

    def test_ooxml_fallback_without_openpyxl(self, monkeypatch: pytest.MonkeyPatch) -> None:
        monkeypatch.setattr(document_extractor_module, "load_workbook", None)
        data = _make_xlsx({"Alpha": [["name", "score"], ["alice", 98]]})

        text = extract_text_from_bytes("book.xlsx", data)

        assert "--- Sheet: Alpha ---" in text
        assert "alice" in text
        assert "98" in text


class TestExtractPptx:
    def test_basic_slides(self) -> None:
        data = _make_pptx([["Slide 1 title", "Slide 1 body"], ["Slide 2 only text"]])
        text = extract_text_from_bytes("deck.pptx", data)
        assert "--- Slide 1 ---" in text
        assert "--- Slide 2 ---" in text
        assert "Slide 1 title" in text
        assert "Slide 2 only text" in text

    def test_ooxml_fallback_without_python_pptx(self, monkeypatch: pytest.MonkeyPatch) -> None:
        monkeypatch.setattr(document_extractor_module, "PptxPresentation", None)
        data = _make_pptx([["Fallback slide", "第二行"]])

        text = extract_text_from_bytes("deck.pptx", data)

        assert "--- Slide 1 ---" in text
        assert "Fallback slide" in text
        assert "第二行" in text


class TestExtractTextLike:
    def test_plain_txt(self) -> None:
        text = extract_text_from_bytes("note.txt", "hello world\nline two".encode("utf-8"))
        assert "hello world" in text
        assert "line two" in text

    def test_python_source(self) -> None:
        src = b"def greet(name: str) -> str:\n    return f'hi {name}'\n"
        text = extract_text_from_bytes("greet.py", src)
        assert "def greet" in text

    def test_json(self) -> None:
        text = extract_text_from_bytes("data.json", b'{"x": 42, "y": "ok"}')
        assert '"x": 42' in text

    def test_csv(self) -> None:
        text = extract_text_from_bytes("table.csv", b"a,b,c\n1,2,3\n")
        assert "a,b,c" in text
        assert "1,2,3" in text

    def test_markdown(self) -> None:
        text = extract_text_from_bytes("doc.md", b"# Heading\n\nBody.\n")
        assert "# Heading" in text

    def test_utf8_with_bom(self) -> None:
        # The candidate chain has utf-8 before utf-8-sig (same order as KB
        # pipeline), so BOM-prefixed bytes decode as utf-8 and the BOM is
        # retained. Mirror that behavior here.
        text = extract_text_from_bytes("note.txt", b"\xef\xbb\xbfhello")
        assert "hello" in text

    def test_gbk_fallback(self) -> None:
        # "你好" in GBK
        text = extract_text_from_bytes("note.txt", "你好".encode("gbk"))
        assert text == "你好"

    def test_svg(self) -> None:
        svg = (
            b'<?xml version="1.0"?>'
            b'<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">'
            b'<circle cx="50" cy="50" r="40" fill="red"/>'
            b'<text x="50" y="55">Hello</text>'
            b"</svg>"
        )
        text = extract_text_from_bytes("logo.svg", svg)
        assert "<svg" in text
        assert "<circle" in text
        assert "Hello" in text


class TestExtractPdf:
    def test_minimal_pdf(self) -> None:
        # Build a minimal valid PDF via pymupdf (dependency already in project).
        pytest.importorskip("fitz")
        import fitz  # noqa: WPS433

        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Hello PDF world")
        buf = io.BytesIO()
        doc.save(buf)
        doc.close()
        data = buf.getvalue()

        text = extract_text_from_bytes("sample.pdf", data)
        assert "Hello PDF world" in text
        assert "--- Page 1 ---" in text


# ---------------------------------------------------------------------------
# Failure modes
# ---------------------------------------------------------------------------


class TestFailureModes:
    def test_unsupported_extension(self) -> None:
        with pytest.raises(UnsupportedDocumentError):
            extract_text_from_bytes("foo.zip", b"\x00\x00")

    def test_empty_bytes(self) -> None:
        with pytest.raises(EmptyDocumentError):
            extract_text_from_bytes("foo.docx", b"")

    def test_too_large(self) -> None:
        fake = b"PK\x03\x04" + b"\x00" * (MAX_DOC_BYTES + 1)
        with pytest.raises(DocumentTooLargeError):
            extract_text_from_bytes("foo.docx", fake)

    def test_pdf_magic_mismatch(self) -> None:
        with pytest.raises(CorruptDocumentError):
            extract_text_from_bytes("foo.pdf", b"this is not a pdf")

    def test_ooxml_magic_mismatch(self) -> None:
        with pytest.raises(CorruptDocumentError):
            extract_text_from_bytes("foo.docx", b"not an office file")

    def test_corrupt_docx(self) -> None:
        # OOXML header but garbage body
        with pytest.raises(CorruptDocumentError):
            extract_text_from_bytes("foo.docx", b"PK\x03\x04" + b"\x00" * 512)

    def test_empty_docx_no_text(self) -> None:
        data = _make_docx([])  # no paragraphs
        with pytest.raises(EmptyDocumentError):
            extract_text_from_bytes("foo.docx", data)


# ---------------------------------------------------------------------------
# Truncation
# ---------------------------------------------------------------------------


class TestTruncation:
    def test_long_docx_is_truncated(self) -> None:
        # single paragraph of 250k chars → well over the 200k per-doc cap
        long_text = "a" * 250_000
        data = _make_docx([long_text])
        text = extract_text_from_bytes("big.docx", data)
        assert len(text) <= MAX_EXTRACTED_CHARS_PER_DOC + 200  # allow notice suffix
        assert "truncated" in text


# ---------------------------------------------------------------------------
# extract_documents_from_records
# ---------------------------------------------------------------------------


class TestExtractDocumentsFromRecords:
    def test_mixed_image_and_doc(self) -> None:
        docx_bytes = _make_docx(["hello there"])
        docx_b64 = base64.b64encode(docx_bytes).decode()
        image_b64 = base64.b64encode(b"\x89PNG\r\n\x1a\n").decode()

        records = [
            {
                "type": "image",
                "filename": "pic.png",
                "base64": image_b64,
                "mime_type": "image/png",
                "url": "",
            },
            {
                "type": "file",
                "filename": "note.docx",
                "base64": docx_b64,
                "mime_type": "",
                "url": "",
            },
        ]

        doc_texts, updated = extract_documents_from_records(records)

        assert len(doc_texts) == 1
        assert "[File: note.docx]" in doc_texts[0]
        assert "hello there" in doc_texts[0]

        # image record untouched
        assert updated[0]["base64"] == image_b64
        # doc record base64 cleared, extracted_chars set
        assert updated[1]["base64"] == ""
        assert updated[1]["extracted_chars"] > 0

    def test_unsupported_record_is_passthrough(self) -> None:
        records = [
            {"type": "file", "filename": "foo.zip", "base64": "AAAA", "mime_type": "", "url": ""}
        ]
        doc_texts, updated = extract_documents_from_records(records)
        assert doc_texts == []
        assert updated[0]["base64"] == "AAAA"  # untouched — not a doc extension

    def test_failed_extraction_emits_error_marker(self) -> None:
        records = [
            {
                "type": "file",
                "filename": "bad.pdf",
                "base64": base64.b64encode(b"not a pdf").decode(),
                "mime_type": "",
                "url": "",
            }
        ]
        doc_texts, updated = extract_documents_from_records(records)
        assert len(doc_texts) == 1
        assert "bad.pdf" in doc_texts[0]
        assert "could not be read" in doc_texts[0]
        assert updated[0]["base64"] == ""  # stripped even on failure

    def test_invalid_base64_emits_error_marker(self) -> None:
        records = [
            {
                "type": "file",
                "filename": "bad.docx",
                "base64": "!!!not base64!!!",
                "mime_type": "",
                "url": "",
            }
        ]
        doc_texts, updated = extract_documents_from_records(records)
        # invalid base64 with validate=False may silently decode or emit error — both
        # paths end up as an error marker since resulting bytes won't pass magic check
        assert len(doc_texts) == 1
        assert "bad.docx" in doc_texts[0]
