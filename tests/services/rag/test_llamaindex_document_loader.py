"""Tests for LlamaIndex document loading."""

from __future__ import annotations

import asyncio
import io
from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
import pytest


def _make_docx(paragraphs: list[str]) -> bytes:
    doc = DocxDocument()
    for paragraph in paragraphs:
        doc.add_paragraph(paragraph)
    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def _make_xlsx(rows: list[list[object]]) -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    for row in rows:
        sheet.append(row)
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _make_pptx(texts: list[str]) -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    for index, text in enumerate(texts):
        text_box = slide.shapes.add_textbox(
            Inches(1),
            Inches(1 + index * 0.5),
            Inches(6),
            Inches(0.5),
        )
        text_box.text_frame.text = text
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_loader_extracts_chat_supported_office_files(tmp_path: Path) -> None:
    pytest.importorskip("llama_index.core")
    from deeptutor.services.rag.pipelines.llamaindex.document_loader import (
        LlamaIndexDocumentLoader,
    )

    docx_path = tmp_path / "notes.docx"
    docx_path.write_bytes(_make_docx(["Docx paragraph"]))
    xlsx_path = tmp_path / "book.xlsx"
    xlsx_path.write_bytes(_make_xlsx([["cell-a", 42]]))
    pptx_path = tmp_path / "slides.pptx"
    pptx_path.write_bytes(_make_pptx(["Slide title", "Slide body"]))

    documents = asyncio.run(
        LlamaIndexDocumentLoader().load([str(docx_path), str(xlsx_path), str(pptx_path)])
    )

    assert {doc.metadata["file_name"] for doc in documents} == {
        "notes.docx",
        "book.xlsx",
        "slides.pptx",
    }
    all_text = "\n".join(doc.text for doc in documents)
    assert "Docx paragraph" in all_text
    assert "cell-a" in all_text
    assert "42" in all_text
    assert "Slide title" in all_text
    assert "Slide body" in all_text
