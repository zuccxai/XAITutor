"""Document text extraction for chat attachments.

Bytes-in, text-out. Used by the chat turn runtime to inline the text of
user-dropped files into the ``effective_user_message`` sent to the LLM.

Two format families:
  * **Binary Office** (.pdf / .docx / .xlsx / .pptx) — parsed with pymupdf /
    python-docx / openpyxl / python-pptx.
  * **Text-like** (plain text, Markdown, source code, JSON, XML, CSV, …) —
    the extension set is imported from ``FileTypeRouter.TEXT_EXTENSIONS`` so
    the chat composer accepts every format the knowledge-base pipeline
    already ingests. Decoded with the same multi-encoding fallback chain.

Design mirrors ``nanobot/nanobot/utils/document.py`` but works on bytes
instead of file paths so the server never touches disk.
"""

from __future__ import annotations

import base64
from collections.abc import Iterable
import io
import logging
from pathlib import Path, PurePosixPath
import re
from typing import Any
import zipfile

from defusedxml import ElementTree as DefusedElementTree
from defusedxml.common import DefusedXmlException

from deeptutor.services.rag.file_routing import FileTypeRouter

try:
    import fitz  # pymupdf
except ImportError:  # pragma: no cover
    fitz = None  # type: ignore[assignment]

try:
    from pypdf import PdfReader
    from pypdf.errors import FileNotDecryptedError as _PypdfNotDecryptedError
except ImportError:  # pragma: no cover
    PdfReader = None  # type: ignore[assignment]
    _PypdfNotDecryptedError = Exception  # type: ignore[assignment,misc]

try:
    from docx import Document as DocxDocument
except ImportError:  # pragma: no cover
    DocxDocument = None  # type: ignore[assignment]

try:
    from openpyxl import load_workbook
except ImportError:  # pragma: no cover
    load_workbook = None  # type: ignore[assignment]

try:
    from pptx import Presentation as PptxPresentation
except ImportError:  # pragma: no cover
    PptxPresentation = None  # type: ignore[assignment]


logger = logging.getLogger(__name__)


_OFFICE_EXTENSIONS: frozenset[str] = frozenset(FileTypeRouter.PARSER_EXTENSIONS)
# Text-like formats are sourced from the KB file router so chat and KB stay
# in sync. Adding a new code / config extension in one place propagates here.
TEXT_LIKE_EXTENSIONS: frozenset[str] = frozenset(FileTypeRouter.TEXT_EXTENSIONS)
SUPPORTED_DOC_EXTENSIONS: frozenset[str] = _OFFICE_EXTENSIONS | TEXT_LIKE_EXTENSIONS

MAX_DOC_BYTES = 10 * 1024 * 1024
MAX_TOTAL_DOC_BYTES = 25 * 1024 * 1024
MAX_EXTRACTED_CHARS_PER_DOC = 200_000
MAX_EXTRACTED_CHARS_TOTAL = 150_000

_PDF_MAGIC = b"%PDF-"
_OOXML_MAGIC = b"PK\x03\x04"


class DocumentExtractionError(Exception):
    """Base class for extraction failures. ``str(exc)`` is user-friendly."""

    def __init__(self, message: str, filename: str = "") -> None:
        super().__init__(message)
        self.filename = filename


class UnsupportedDocumentError(DocumentExtractionError):
    pass


class CorruptDocumentError(DocumentExtractionError):
    pass


class EmptyDocumentError(DocumentExtractionError):
    pass


class DocumentTooLargeError(DocumentExtractionError):
    pass


def is_document_extension(filename: str) -> bool:
    return _ext(filename) in SUPPORTED_DOC_EXTENSIONS


def _ext(filename: str) -> str:
    return PurePosixPath(filename or "").suffix.lower()


def _truncate(text: str, max_length: int) -> str:
    if len(text) <= max_length:
        return text
    return text[:max_length] + f"... (truncated, {len(text)} chars total)"


def _check_magic(ext: str, data: bytes, filename: str) -> None:
    """Validate file header to catch extension spoofing.

    Only binary formats have well-known magic prefixes. Text-like extensions
    (code, markup, config, …) are decoded directly; a mislabeled binary blob
    either decodes as garbage or fails at decode time, which is fine.
    """
    if ext == ".pdf":
        if not data.startswith(_PDF_MAGIC):
            raise CorruptDocumentError(
                f"{filename} does not look like a PDF (bad header)", filename=filename
            )
    elif ext in {".docx", ".xlsx", ".pptx"}:
        if not data.startswith(_OOXML_MAGIC):
            raise CorruptDocumentError(
                f"{filename} does not look like a valid Office file (bad header)",
                filename=filename,
            )


def extract_text_from_bytes(
    filename: str,
    data: bytes,
    *,
    max_bytes: int | None = MAX_DOC_BYTES,
    max_chars: int | None = MAX_EXTRACTED_CHARS_PER_DOC,
) -> str:
    """Extract text from a single document's raw bytes.

    Raises a ``DocumentExtractionError`` subclass on failure. Successful
    output is truncated to ``max_chars`` with a notice when ``max_chars`` is
    not ``None``. ``max_bytes`` is configurable so the KB indexer can reuse
    the same parsers with its larger upload policy while chat keeps the
    stricter per-turn limit.
    """
    if not data:
        raise EmptyDocumentError(f"{filename} is empty", filename=filename)
    if max_bytes is not None and len(data) > max_bytes:
        raise DocumentTooLargeError(
            f"{filename} exceeds the {max_bytes // (1024 * 1024)} MB per-file limit",
            filename=filename,
        )

    ext = _ext(filename)
    if ext not in SUPPORTED_DOC_EXTENSIONS:
        raise UnsupportedDocumentError(
            f"{filename} has unsupported extension '{ext}'", filename=filename
        )

    _check_magic(ext, data, filename)

    if ext == ".pdf":
        text = _extract_pdf(data, filename)
    elif ext == ".docx":
        text = _extract_docx(data, filename)
    elif ext == ".xlsx":
        text = _extract_xlsx(data, filename)
    elif ext == ".pptx":
        text = _extract_pptx(data, filename)
    elif ext in TEXT_LIKE_EXTENSIONS:
        text = _extract_text_like(data, filename)
    else:  # pragma: no cover - guarded above
        raise UnsupportedDocumentError(f"{filename}: unreachable", filename=filename)

    if not text.strip():
        raise EmptyDocumentError(f"{filename}: no extractable text", filename=filename)
    return _truncate(text, max_chars) if max_chars is not None else text


def extract_text_from_path(
    file_path: str | Path,
    *,
    max_bytes: int | None = MAX_DOC_BYTES,
    max_chars: int | None = MAX_EXTRACTED_CHARS_PER_DOC,
) -> str:
    """Extract text from a file path using the same bytes-based parsers."""
    path = Path(file_path)
    return extract_text_from_bytes(
        path.name,
        path.read_bytes(),
        max_bytes=max_bytes,
        max_chars=max_chars,
    )


def _extract_pdf(data: bytes, filename: str) -> str:
    if fitz is not None:
        try:
            with fitz.open(stream=data, filetype="pdf") as doc:
                if doc.is_encrypted and not doc.authenticate(""):
                    raise CorruptDocumentError(
                        f"{filename} is encrypted and cannot be read", filename=filename
                    )
                pages = [
                    f"--- Page {i} ---\n{page.get_text() or ''}" for i, page in enumerate(doc, 1)
                ]
            return "\n\n".join(pages)
        except CorruptDocumentError:
            raise
        except Exception as exc:
            logger.warning("pymupdf failed on %s: %s — falling back to pypdf", filename, exc)

    if PdfReader is None:
        raise CorruptDocumentError(
            f"{filename}: no PDF reader available (install pymupdf or pypdf)",
            filename=filename,
        )
    try:
        reader = PdfReader(io.BytesIO(data))
        if getattr(reader, "is_encrypted", False):
            raise CorruptDocumentError(
                f"{filename} is encrypted and cannot be read", filename=filename
            )
        pages = [
            f"--- Page {i} ---\n{page.extract_text() or ''}"
            for i, page in enumerate(reader.pages, 1)
        ]
        return "\n\n".join(pages)
    except CorruptDocumentError:
        raise
    except _PypdfNotDecryptedError as exc:
        raise CorruptDocumentError(
            f"{filename} is encrypted and cannot be read", filename=filename
        ) from exc
    except Exception as exc:
        raise CorruptDocumentError(
            f"{filename}: failed to read PDF ({exc})", filename=filename
        ) from exc


def _extract_docx(data: bytes, filename: str) -> str:
    primary_error: Exception | None = None
    primary_text = ""
    if DocxDocument is not None:
        try:
            doc = DocxDocument(io.BytesIO(data))
            paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
            primary_text = "\n\n".join(paragraphs)
        except Exception as exc:
            primary_error = exc
            logger.info("python-docx failed on %s; falling back to raw OOXML: %s", filename, exc)

    fallback = _extract_docx_ooxml(data, filename)
    if fallback.strip() and (not primary_text.strip() or len(fallback) > len(primary_text) * 1.2):
        return fallback
    if primary_text.strip():
        return primary_text

    if DocxDocument is None:
        raise CorruptDocumentError(
            f"{filename}: python-docx not installed and OOXML fallback found no text",
            filename=filename,
        )
    if primary_error is not None:
        raise CorruptDocumentError(
            f"{filename}: failed to open DOCX ({primary_error})", filename=filename
        ) from primary_error
    return ""


def _extract_xlsx(data: bytes, filename: str) -> str:
    if load_workbook is None:
        return _extract_xlsx_ooxml(data, filename)
    try:
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as exc:
        logger.info("openpyxl failed on %s; falling back to raw OOXML: %s", filename, exc)
        fallback = _extract_xlsx_ooxml(data, filename)
        if fallback.strip():
            return fallback
        raise CorruptDocumentError(
            f"{filename}: failed to open XLSX ({exc})", filename=filename
        ) from exc
    try:
        sheets: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                sheets.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
        return "\n\n".join(sheets)
    finally:
        wb.close()


def _extract_pptx(data: bytes, filename: str) -> str:
    if PptxPresentation is None:
        return _extract_pptx_ooxml(data, filename)
    try:
        prs = PptxPresentation(io.BytesIO(data))
    except Exception as exc:
        logger.info("python-pptx failed on %s; falling back to raw OOXML: %s", filename, exc)
        fallback = _extract_pptx_ooxml(data, filename)
        if fallback.strip():
            return fallback
        raise CorruptDocumentError(
            f"{filename}: failed to open PPTX ({exc})", filename=filename
        ) from exc
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            _collect_pptx_shape_text(shape, slide_text)
        if slide_text:
            slides.append(f"--- Slide {i} ---\n" + "\n".join(slide_text))
    return "\n\n".join(slides)


def _extract_text_like(data: bytes, filename: str) -> str:
    """Decode a plain-text / code / config / markup file.

    Uses the same encoding fallback chain as the KB pipeline
    (``FileTypeRouter.decode_bytes``) so a GBK-encoded Python file or a
    UTF-8-BOM Markdown works the same way in both places.
    """
    try:
        return FileTypeRouter.decode_bytes(data)
    except Exception as exc:  # pragma: no cover - decode_bytes never raises
        raise CorruptDocumentError(
            f"{filename}: failed to decode text ({exc})", filename=filename
        ) from exc


def _open_ooxml(data: bytes, filename: str) -> zipfile.ZipFile:
    try:
        return zipfile.ZipFile(io.BytesIO(data))
    except zipfile.BadZipFile as exc:
        raise CorruptDocumentError(
            f"{filename}: failed to open Office ZIP package ({exc})", filename=filename
        ) from exc


def _local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1] if "}" in tag else tag


def _parse_xml_member(zf: zipfile.ZipFile, member: str, filename: str) -> Any | None:
    try:
        raw = zf.read(member)
    except KeyError:
        return None
    try:
        return DefusedElementTree.fromstring(raw)
    except (DefusedElementTree.ParseError, DefusedXmlException) as exc:
        raise CorruptDocumentError(
            f"{filename}: failed to parse {member} ({exc})", filename=filename
        ) from exc


def _collect_ooxml_text(node: Any) -> str:
    parts: list[str] = []
    for child in node.iter():
        name = _local_name(child.tag)
        if name == "t" and child.text:
            parts.append(child.text)
        elif name == "tab":
            parts.append("\t")
        elif name in {"br", "cr"}:
            parts.append("\n")
    return "".join(parts).strip()


def _extract_paragraph_text(root: Any) -> list[str]:
    paragraphs: list[str] = []
    for node in root.iter():
        if _local_name(node.tag) != "p":
            continue
        text = _collect_ooxml_text(node)
        if text:
            paragraphs.append(text)
    if paragraphs:
        return paragraphs
    text = _collect_ooxml_text(root)
    return [text] if text else []


def _extract_docx_ooxml(data: bytes, filename: str) -> str:
    with _open_ooxml(data, filename) as zf:
        names = zf.namelist()
        content_members = ["word/document.xml"]
        content_members.extend(
            sorted(
                name
                for name in names
                if re.match(r"word/(header|footer|footnotes|endnotes|comments)\d*\.xml$", name)
            )
        )

        chunks: list[str] = []
        for member in content_members:
            root = _parse_xml_member(zf, member, filename)
            if root is None:
                continue
            chunks.extend(_extract_paragraph_text(root))
        return "\n\n".join(chunks)


def _xlsx_shared_strings(zf: zipfile.ZipFile, filename: str) -> list[str]:
    root = _parse_xml_member(zf, "xl/sharedStrings.xml", filename)
    if root is None:
        return []
    strings: list[str] = []
    for node in root:
        if _local_name(node.tag) != "si":
            continue
        strings.append(_collect_ooxml_text(node))
    return strings


def _xlsx_sheet_names(zf: zipfile.ZipFile, filename: str) -> dict[str, str]:
    root = _parse_xml_member(zf, "xl/workbook.xml", filename)
    if root is None:
        return {}
    out: dict[str, str] = {}
    index = 1
    for node in root.iter():
        if _local_name(node.tag) != "sheet":
            continue
        sheet_name = node.attrib.get("name") or f"sheet{index}"
        sheet_id = node.attrib.get("sheetId") or str(index)
        out[f"xl/worksheets/sheet{sheet_id}.xml"] = sheet_name
        index += 1
    return out


def _xlsx_cell_text(cell: Any, shared_strings: list[str]) -> str:
    cell_type = cell.attrib.get("t", "")
    if cell_type == "inlineStr":
        return _collect_ooxml_text(cell)

    value = ""
    for child in cell:
        if _local_name(child.tag) == "v":
            value = child.text or ""
            break

    if cell_type == "s":
        try:
            return shared_strings[int(value)]
        except (ValueError, IndexError):
            return value
    return value


def _extract_xlsx_ooxml(data: bytes, filename: str) -> str:
    with _open_ooxml(data, filename) as zf:
        shared_strings = _xlsx_shared_strings(zf, filename)
        sheet_names = _xlsx_sheet_names(zf, filename)
        sheet_members = sorted(
            (name for name in zf.namelist() if re.match(r"xl/worksheets/sheet\d+\.xml$", name)),
            key=lambda name: [
                int(part) if part.isdigit() else part for part in re.split(r"(\d+)", name)
            ],
        )

        sheets: list[str] = []
        for index, member in enumerate(sheet_members, 1):
            root = _parse_xml_member(zf, member, filename)
            if root is None:
                continue
            rows: list[str] = []
            for row in root.iter():
                if _local_name(row.tag) != "row":
                    continue
                cells = [
                    _xlsx_cell_text(cell, shared_strings)
                    for cell in row
                    if _local_name(cell.tag) == "c"
                ]
                row_text = "\t".join(cells)
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                sheet_name = sheet_names.get(member, f"sheet{index}")
                sheets.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
        return "\n\n".join(sheets)


def _extract_pptx_ooxml(data: bytes, filename: str) -> str:
    with _open_ooxml(data, filename) as zf:
        slide_members = sorted(
            (name for name in zf.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", name)),
            key=lambda name: [
                int(part) if part.isdigit() else part for part in re.split(r"(\d+)", name)
            ],
        )
        slides: list[str] = []
        for index, member in enumerate(slide_members, 1):
            root = _parse_xml_member(zf, member, filename)
            if root is None:
                continue
            paragraphs = _extract_paragraph_text(root)
            if paragraphs:
                slides.append(f"--- Slide {index} ---\n" + "\n".join(paragraphs))
        return "\n\n".join(slides)


def _collect_pptx_shape_text(shape, out: list[str]) -> None:
    """Recurse into groups + tables, same semantics as nanobot's version."""
    sub_shapes = getattr(shape, "shapes", None)
    if sub_shapes is not None:
        for sub in sub_shapes:
            _collect_pptx_shape_text(sub, out)
        return

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            line = "\t".join(cell for cell in cells if cell)
            if line:
                out.append(line)
        return

    text = getattr(shape, "text", "")
    if text:
        out.append(text)


def extract_documents_from_records(
    records: Iterable[dict],
) -> tuple[list[str], list[dict]]:
    """Process a list of attachment records from the WS payload.

    Parameters
    ----------
    records:
        Raw attachment records as parsed by the turn runtime
        (``{"type", "url", "base64", "filename", "mime_type"}``).

    Returns
    -------
    (doc_texts, updated_records)
        ``doc_texts`` is a list of strings formatted as
        ``"[File: <name>]\\n<text>"`` (one per processed or skipped doc).
        ``updated_records`` is the input list with the ``base64`` field
        cleared on successfully-extracted docs (to save DB space), an
        ``extracted_chars`` field added, and the extracted plain text
        stored under ``extracted_text`` so the chat UI can preview office
        documents without re-running the parser. Image / non-document
        records are returned unchanged.
    """
    doc_texts: list[str] = []
    updated: list[dict] = []
    total_bytes = 0
    total_chars = 0
    over_quota = False

    for raw in records:
        record = dict(raw)
        filename = str(record.get("filename") or "")
        if not is_document_extension(filename):
            updated.append(record)
            continue

        b64 = record.get("base64") or ""
        if not b64:
            updated.append(record)
            continue

        if over_quota:
            doc_texts.append(f"[File: {filename} — skipped: total attachment quota exceeded]")
            record["base64"] = ""
            record["extracted_chars"] = 0
            updated.append(record)
            continue

        try:
            data = base64.b64decode(b64, validate=False)
        except Exception as exc:
            doc_texts.append(f"[File: {filename} — could not be read: invalid base64 ({exc})]")
            record["base64"] = ""
            record["extracted_chars"] = 0
            updated.append(record)
            continue

        if total_bytes + len(data) > MAX_TOTAL_DOC_BYTES:
            over_quota = True
            doc_texts.append(f"[File: {filename} — skipped: total attachment quota exceeded]")
            record["base64"] = ""
            record["extracted_chars"] = 0
            updated.append(record)
            continue

        total_bytes += len(data)

        try:
            text = extract_text_from_bytes(filename, data)
        except DocumentExtractionError as exc:
            logger.info("Document extraction failed for %s: %s", filename, exc)
            doc_texts.append(f"[File: {filename} — could not be read: {exc}]")
            record["base64"] = ""
            record["extracted_chars"] = 0
            updated.append(record)
            continue

        remaining_budget = MAX_EXTRACTED_CHARS_TOTAL - total_chars
        if remaining_budget <= 0:
            doc_texts.append(f"[File: {filename} — skipped: total extracted-text quota exceeded]")
            record["base64"] = ""
            record["extracted_chars"] = 0
            updated.append(record)
            continue

        if len(text) > remaining_budget:
            text = (
                text[:remaining_budget]
                + f"... (truncated, {len(text)} chars total; turn quota hit)"
            )

        total_chars += len(text)
        doc_texts.append(f"[File: {filename}]\n{text}")
        record["base64"] = ""
        record["extracted_chars"] = len(text)
        record["extracted_text"] = text
        updated.append(record)

    return doc_texts, updated
